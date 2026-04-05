from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# COLORS
BG_DARK = RGBColor(0x1A,0x1A,0x2E)
BG_CARD = RGBColor(0x16,0x21,0x3E)
BG_LIGHTER = RGBColor(0x1F,0x2B,0x4D)
C_BLUE = RGBColor(0x00,0xBF,0xFF)
C_GREEN = RGBColor(0x00,0xE6,0x76)
C_ORANGE = RGBColor(0xFF,0x9F,0x43)
C_RED = RGBColor(0xFF,0x6B,0x6B)
C_PURPLE = RGBColor(0xA2,0x9B,0xFE)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LGRAY = RGBColor(0xCC,0xCC,0xCC)
MGRAY = RGBColor(0x99,0x99,0x99)
DTXT = RGBColor(0xE0,0xE0,0xE0)
YELLOW = RGBColor(0xFF,0xD9,0x3D)
CODE_BG = RGBColor(0x0D,0x11,0x17)
CODE_BD = RGBColor(0x30,0x36,0x3D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, c):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def rect(slide, l, t, w, h, fc, bc=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = bw
    else: s.line.fill.background()
    return s

def tbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def stxt(tf, txt, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn="Segoe UI"):
    tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = fn
    return p

def apara(tf, txt, sz=14, c=DTXT, b=False, al=PP_ALIGN.LEFT, fn="Segoe UI"):
    p = tf.add_paragraph(); p.alignment = al; p.space_before = Pt(3); p.space_after = Pt(3)
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = fn
    return p

def code_blk(slide, l, t, w, h, txt, fs=11):
    s = rect(slide, l, t, w, h, CODE_BG, CODE_BD, Pt(1))
    stxt(s.text_frame, txt, sz=fs, c=C_GREEN, fn="Consolas")
    return s

def title_slide(title, sub, ac=C_BLUE):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG_DARK)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ac; bar.line.fill.background()
    tb = tbox(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    stxt(tb.text_frame, title, sz=40, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    tb2 = tbox(sl, Inches(1.5), Inches(3.6), Inches(10), Inches(1.0))
    stxt(tb2.text_frame, sub, sz=20, c=LGRAY, al=PP_ALIGN.CENTER)
    bar2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = ac; bar2.line.fill.background()
    return sl

def content_slide(title, ac=C_BLUE):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG_DARK)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = ac; bar.line.fill.background()
    tbg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.04), Inches(13.333), Inches(0.75))
    tbg.fill.solid(); tbg.fill.fore_color.rgb = BG_CARD; tbg.line.fill.background()
    tb = tbox(sl, Inches(0.5), Inches(0.1), Inches(12), Inches(0.65))
    stxt(tb.text_frame, title, sz=24, c=ac, b=True)
    return sl

def flow_table(slide, l, t, w, rows, hc=C_BLUE):
    rh = Inches(0.35)
    s = rect(slide, l, t, w, rh, hc)
    stxt(s.text_frame, "  Step  |  Layer / Component  |  Action", sz=11, c=BG_DARK, b=True, fn="Consolas")
    y = t + rh
    for i, (step, layer, action) in enumerate(rows):
        bgc = BG_LIGHTER if i%2==0 else BG_CARD
        rs = rect(slide, l, y, w, rh, bgc)
        lc = DTXT
        if any(k in layer for k in ["Frontend","User","clicks"]): lc = C_BLUE
        elif any(k in layer for k in ["Service","Controller","Check"]): lc = C_GREEN
        elif any(k in layer.upper() for k in ["REPOSITORY","INSERT","SELECT","UPDATE","DELETE"]): lc = C_ORANGE
        stxt(rs.text_frame, f"  {step}  |  {layer}  |  {action}", sz=10, c=lc, fn="Consolas")
        y += rh
    return y

def qa_slide(title, qas, ac=C_RED):
    sl = content_slide(title, ac); y = Inches(0.95)
    for q, a in qas:
        qs = rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.4), ac)
        stxt(qs.text_frame, f"  {q}", sz=12, c=BG_DARK, b=True)
        y += Inches(0.42)
        ans = rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.7), BG_CARD)
        ans.text_frame.word_wrap = True
        stxt(ans.text_frame, f"  {a}", sz=10, c=LGRAY)
        y += Inches(0.75)
    return sl

print("Generating PPT...")

# === SLIDE 1: TITLE ===
sl = title_slide("RevConnect","Full-Stack Social Media Application  |  Complete Technical Deep-Dive")
tb = tbox(sl, Inches(2.5), Inches(4.8), Inches(8), Inches(0.5))
stxt(tb.text_frame, "Spring Boot  |  Angular 19  |  MySQL  |  JWT  |  Hibernate  |  REST API", sz=14, c=C_PURPLE, al=PP_ALIGN.CENTER)

# === SLIDE 2: TOC ===
sl = content_slide("Table of Contents", C_PURPLE)
left_items = ["1. Project Overview & Architecture","2. Tech Stack","3. Project Structure","4. Request Lifecycle","5. Database Schema","6. Security & JWT","7. Registration & Login","8. Post Management","9. Like/Unlike","10. Comments","11. Share/Repost","12. Connections"]
right_items = ["13. Bookmarks & Notifications","14. User Profile","15. Direct Messaging","16. Stories","17. Analytics Dashboard","18. Search & Discovery","19. Media Upload","20. Business/Creator Features","21. Cross-Cutting (AOP)","22. End-to-End Flow","23. API Reference","24. 35 Interview Q&A"]
tl = tbox(sl, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.5)); tf = tl.text_frame
stxt(tf, left_items[0], sz=14, c=C_BLUE)
for it in left_items[1:]: apara(tf, it, sz=14, c=C_BLUE)
tr = tbox(sl, Inches(6.8), Inches(1.1), Inches(5.8), Inches(5.5)); tf2 = tr.text_frame
stxt(tf2, right_items[0], sz=14, c=C_GREEN)
for it in right_items[1:]: apara(tf2, it, sz=14, c=C_GREEN)

# === SLIDE 3: ARCHITECTURE ===
sl = content_slide("3D Architecture Overview", C_BLUE)
lw = Inches(10.5); y = Inches(1.1); g = Inches(0.12)
layers = [
    ("BROWSER / CLIENT LAYER","Angular 19 | Standalone Components | Lazy Routes | HTTP Interceptor | Auth Guard | localStorage JWT",RGBColor(0x0A,0x2A,0x5E),C_BLUE),
    ("SECURITY LAYER","Spring Security + JWT Filter | JwtAuthenticationFilter | BCrypt PasswordEncoder",RGBColor(0x2D,0x13,0x2A),C_RED),
    ("CONTROLLER LAYER  |  @RestController","AuthController | PostController | UserController | InteractionController | +12 more",RGBColor(0x0A,0x3D,0x0A),C_GREEN),
    ("SERVICE LAYER  |  @Service + @Transactional","Business Logic | Validation | Authorization | Notification Triggers | DTO Conversion",RGBColor(0x0A,0x3D,0x0A),C_GREEN),
    ("REPOSITORY LAYER  |  Spring Data JPA + Hibernate","PostRepository | UserRepository | LikeRepository | ConnectionRepository | +8 more",RGBColor(0x3D,0x25,0x0A),C_ORANGE),
    ("MySQL DATABASE  |  17+ Tables | Indexes | Constraints","ddl-auto=update | UNIQUE constraints | Foreign Keys | Denormalized counters",RGBColor(0x3D,0x25,0x0A),C_ORANGE),
]
for name, desc, bgc, bc in layers:
    s = rect(sl, Inches(1.4), y, lw, Inches(0.85), bgc, bc, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    stxt(tf, f"  {name}", sz=13, c=bc, b=True); apara(tf, f"  {desc}", sz=10, c=LGRAY)
    y += Inches(0.85) + g

# === SLIDE 4: TECH STACK ===
sl = content_slide("Tech Stack - WHY Each Choice", C_GREEN)
techs = [
    ("Spring Boot 3.x","Auto-config, embedded Tomcat, production-ready, massive ecosystem",C_GREEN),
    ("Angular 19","Component-based, TypeScript, standalone components, enterprise-grade",C_BLUE),
    ("MySQL","ACID-compliant relational DB, perfect for structured social data with JOINs",C_ORANGE),
    ("Hibernate / JPA","ORM eliminates raw SQL, entity lifecycle, lazy loading, auto-schema",C_GREEN),
    ("JWT (JSON Web Token)","Stateless auth, no server sessions, horizontal scaling ready",C_RED),
    ("Spring Security","Filter chain, BCrypt, CORS/CSRF config, industry standard",C_RED),
    ("Lombok","Eliminates boilerplate: @Builder, @Data, @RequiredArgsConstructor, @Slf4j",C_GREEN),
    ("Spring AOP","Cross-cutting logging via @Aspect without modifying business code",C_PURPLE),
]
for i,(nm,rs,c) in enumerate(techs):
    x = Inches(0.5)+(i%2)*Inches(6.4); y2 = Inches(1.0)+(i//2)*Inches(1.45)
    cd = rect(sl, x, y2, Inches(6.0), Inches(1.3), BG_CARD, c, Pt(1.5))
    tf = cd.text_frame; tf.word_wrap = True; stxt(tf, f"  {nm}", sz=16, c=c, b=True); apara(tf, f"  {rs}", sz=11, c=LGRAY)

# === SLIDE 5: PROJECT STRUCTURE ===
sl = content_slide("Project Structure", C_PURPLE)
be = "backend/src/main/java/.../revconnect/\n  config/    SecurityConfig, WebSocketConfig\n  controller/ Auth, Post, User, Interaction,\n             Message, Notification, Search,\n             Analytics, Media, +7 more\n  service/   Auth, Post, User, Interaction,\n             Connection, Notification, +6 more\n  repository/ Post, User, Like, Comment,\n             Connection, +8 more\n  model/     User, Post, Comment, Like,\n             Bookmark, Connection, Notification,\n             Message, Story, +5 more\n  dto/       request/ + response/ (ApiResponse,\n             PagedResponse, PostResponse...)\n  security/  JwtTokenProvider, JwtAuthFilter\n  aspect/    LoggingAspect (AOP)\n  enums/     PostType, UserType, ConnectionStatus"
fe = "frontend/src/app/\n  core/\n    components/ navbar/, sidebar/\n    guards/    auth.guard.ts\n    interceptors/ auth.interceptor.ts\n    services/  user, post, auth, message,\n               notification, story, search +5\n  features/\n    auth/      login/, register/, forgot-pw/\n    feed/      feed-page/ (main timeline)\n    explore/   explore-page/ (discover)\n    profile/   profile-page/, edit-profile/\n    messages/  messages-page/ (DMs)\n    notifications/ notifications-page/\n    stories/   stories-feed/, create-story/\n    analytics/ analytics-dashboard/\n    bookmarks/ bookmarks-page/\n    settings/  settings-page/"
code_blk(sl, Inches(0.3), Inches(1.0), Inches(6.3), Inches(5.8), be, 9)
code_blk(sl, Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.8), fe, 9)

# === SLIDE 6: REQUEST LIFECYCLE ===
sl = content_slide("Spring Boot Request Lifecycle", C_GREEN)
steps = [
    ("1. HTTP Request","Browser sends request with Authorization: Bearer <JWT>",C_BLUE),
    ("2. JWT Filter","JwtAuthenticationFilter extracts & validates token",C_RED),
    ("3. SecurityContext","User loaded into SecurityContextHolder (thread-local)",C_RED),
    ("4. Controller","@RestController receives request, validates, delegates to Service",C_GREEN),
    ("5. Service","@Service executes business logic, @Transactional manages DB tx",C_GREEN),
    ("6. Repository","Spring Data JPA converts method name to SQL via Hibernate",C_ORANGE),
    ("7. Database","MySQL executes query, returns result set",C_ORANGE),
    ("8. Response","Entity -> DTO (Mapper) -> ApiResponse<T> -> JSON -> HTTP 200",C_GREEN),
]
y = Inches(1.1)
for st,ds,c in steps:
    s = rect(sl, Inches(1.0), y, Inches(11.3), Inches(0.65), BG_CARD, c, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {st}:  "; r1.font.size = Pt(14); r1.font.color.rgb = c; r1.font.bold = True; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = ds; r2.font.size = Pt(12); r2.font.color.rgb = LGRAY; r2.font.name = "Segoe UI"
    y += Inches(0.75)

# === SLIDE 7: DB SCHEMA ===
sl = content_slide("Database Schema - Entity Relationships", C_ORANGE)
ents = [
    ("User","id, username, email, password,\nbio, profilePicture, coverPhoto,\nuserType, privacy, isVerified",Inches(0.3),Inches(1.1)),
    ("Post","id, content, postType,\nuser_id(FK), mediaUrls,\nlikeCount, commentCount, shareCount",Inches(4.6),Inches(1.1)),
    ("Comment","id, content, user_id(FK),\npost_id(FK), parent_id(FK),\nlikeCount, replyCount",Inches(9.0),Inches(1.1)),
    ("Like","id, user_id(FK), post_id(FK)\nUNIQUE(user, post)",Inches(0.3),Inches(3.6)),
    ("Connection","id, follower_id(FK),\nfollowing_id(FK), status\nUNIQUE(follower, following)",Inches(4.6),Inches(3.6)),
    ("Notification","id, user_id, actor_id,\ntype, message,\nreferenceId, isRead",Inches(9.0),Inches(3.6)),
    ("Story","id, user_id(FK), mediaUrl,\ncaption, expiresAt, isHighlight",Inches(0.3),Inches(5.6)),
    ("Message","id, sender_id, receiver_id,\ncontent, isRead, isDeleted",Inches(4.6),Inches(5.6)),
    ("Bookmark","id, user_id(FK), post_id(FK)\nUNIQUE(user, post)",Inches(9.0),Inches(5.6)),
]
for nm,flds,x,y in ents:
    cd = rect(sl, x, y, Inches(4.0), Inches(1.4), BG_CARD, C_ORANGE, Pt(1.5))
    tf = cd.text_frame; tf.word_wrap = True; stxt(tf, f"  {nm}", sz=14, c=C_ORANGE, b=True); apara(tf, f"  {flds}", sz=9, c=LGRAY, fn="Consolas")

# === SLIDE 8: SECURITY & JWT ===
sl = content_slide("Security & JWT Authentication", C_RED)
jwt = "JWT AUTHENTICATION FLOW:\n\n1. POST /api/auth/login { username, password }\n2. AuthenticationManager.authenticate()\n   -> BCryptPasswordEncoder.matches(raw, hash)\n3. JwtTokenProvider.generateToken(user)\n   -> Claims: { sub: userId, exp: now+24h }\n   -> Sign with HMAC-SHA256\n4. Return: { token: \"eyJhbGci...\" }\n\nON EVERY REQUEST:\n5. Angular interceptor: Authorization: Bearer <token>\n6. JwtAuthenticationFilter:\n   -> Extract token -> Validate -> Set SecurityContext\n7. Controller: authService.getCurrentUser()"
sec = "SECURITY CONFIG:\n\nSecurityFilterChain:\n  .csrf().disable()       // JWT, not cookies\n  .cors().enable()        // Allow Angular\n  .sessionManagement()\n    .STATELESS            // No server sessions\n  .authorizeHttpRequests()\n    .requestMatchers(\n      \"/api/auth/**\",\n      \"/swagger-ui/**\"\n    ).permitAll()         // Public\n    .anyRequest()\n      .authenticated()    // JWT required\n  .addFilterBefore(\n    jwtFilter,\n    UsernamePasswordAuth..)"
code_blk(sl, Inches(0.4), Inches(1.0), Inches(6.3), Inches(5.8), jwt, 11)
code_blk(sl, Inches(6.9), Inches(1.0), Inches(6.1), Inches(5.8), sec, 11)

# === SLIDE 9: AUTH FLOWS ===
sl = content_slide("Feature: Registration & Login", C_BLUE)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User fills form, clicks Register","username, email, password, name"),
    ("2","POST /api/auth/register","AuthController.register()"),
    ("3","AuthService.register()","Check username/email exists -> DuplicateResourceException"),
    ("4","BCrypt.encode(password)","Hash password with random salt"),
    ("5","userRepository.save(user)","INSERT INTO users"),
    ("6","Generate OTP, send email","6-digit code for verification"),
    ("7","Return ApiResponse.success()","{ success: true }"),
], C_BLUE)
tb = tbox(sl, Inches(0.4), Inches(4.0), Inches(12), Inches(0.4)); stxt(tb.text_frame, "LOGIN FLOW:", sz=16, c=C_GREEN, b=True)
flow_table(sl, Inches(0.4), Inches(4.4), Inches(12.5), [
    ("1","User clicks Login","username + password"),
    ("2","POST /api/auth/login","AuthController.login()"),
    ("3","AuthenticationManager.authenticate()","BCrypt.matches(raw, stored)"),
    ("4","JwtTokenProvider.generateToken()","{ sub: userId, exp: +24h }"),
    ("5","localStorage.setItem(token)","Token stored client-side"),
    ("6","router.navigate(['/feed'])","Redirect to feed"),
], C_GREEN)

# === SLIDE 10: POST MANAGEMENT ===
sl = content_slide("Feature: Post Management (CRUD)", C_GREEN)
tb = tbox(sl, Inches(0.4), Inches(1.0), Inches(12), Inches(0.4))
stxt(tb.text_frame, "Post Types: TEXT | IMAGE | VIDEO | PROMOTIONAL | ANNOUNCEMENT | UPDATE | REPOST", sz=13, c=C_ORANGE, b=True)
flow_table(sl, Inches(0.4), Inches(1.5), Inches(12.5), [
    ("1","User types content, clicks Post","feed-page.ts -> createPost()"),
    ("2","POST /api/posts {content, postType}","postService.createPost()"),
    ("3","PostService.createPost()","Get user, build Post entity"),
    ("4","parseMetadata(content)","Extract CTA/Tags/Promo markers"),
    ("5","postRepository.save(post)","INSERT INTO posts"),
    ("6","hashtagService.processHashtags()","Extract #tags, update table"),
    ("7","Return PostResponse","DTO with author + interaction flags"),
], C_GREEN)
tb2 = tbox(sl, Inches(0.4), Inches(4.5), Inches(12), Inches(2.5)); tf = tb2.text_frame
stxt(tf, "Additional Features:", sz=15, c=WHITE, b=True)
for t in ["Edit: PUT /api/posts/{id} -> Ownership check -> Update","Delete: DELETE -> Cascade delete comments, likes, bookmarks","Pin: PATCH /api/posts/{id}/pin -> Toggle pinned (shows at top)","Schedule: POST /api/posts/schedule -> ScheduledExecutorService (in-memory)","Feed: GET /api/posts/personalized -> Posts from followed + self, sorted by date"]:
    apara(tf, f"  {t}", sz=12, c=LGRAY)

# === SLIDE 11: LIKE SYSTEM ===
sl = content_slide("Feature: Like / Unlike System", C_RED)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User clicks Heart icon","toggleLike(post)"),
    ("2","POST /api/posts/{id}/like","interactionService.likePost()"),
    ("3","InteractionService.likePost()","@Transactional"),
    ("4","Check: existsByUserAndPost()","Already liked? -> BadRequestException"),
    ("5","Like.builder().user().post().build()","Create Like entity"),
    ("6","likeRepository.save(like)","INSERT INTO likes"),
    ("7","post.setLikeCount(+1)","Increment denormalized counter"),
    ("8","notificationService.notifyLike()","INSERT INTO notifications"),
    ("9","UI: heart fills red, count+1","Instant feedback"),
], C_RED)
cd = rect(sl, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.0), BG_CARD, C_RED, Pt(2))
tf = cd.text_frame; tf.word_wrap = True
stxt(tf, "  DESIGN: Denormalized Counters", sz=14, c=C_RED, b=True)
apara(tf, "  Instead of COUNT(*) FROM likes on every feed load, store likeCount on Post directly.", sz=12, c=LGRAY)
apara(tf, "  Trade-off: Complex writes but O(1) reads. UNIQUE(user_id, post_id) prevents duplicates.", sz=12, c=LGRAY)
apara(tf, "  Unlike: DELETE FROM likes + post.likeCount-- in same @Transactional method.", sz=12, c=LGRAY)

# Save part 1 marker - continues in next edit
# === SLIDE 12: COMMENTS ===
sl = content_slide("Feature: Comment System (Threaded)", C_GREEN)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User types comment, clicks Submit","submitComment(post)"),
    ("2","POST /api/posts/{id}/comments","{content, parentId?}"),
    ("3","InteractionService.addComment()","@Transactional"),
    ("4","Build Comment entity","user, post, content, parent"),
    ("5","commentRepository.save(comment)","INSERT INTO comments"),
    ("6","post.commentCount++","Increment counter"),
    ("7","If reply: parent.replyCount++","Increment parent reply count"),
    ("8","notificationService.notifyComment()","Notify post owner"),
], C_GREEN)
code_blk(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.5),
"NESTED COMMENTS (Self-Referencing):\n\nComment Table:\n  id | content      | post_id | parent_id | reply_count\n  1  | \"Great post\" | 42      | NULL      | 2          <- Top-level\n  2  | \"Thanks!\"    | 42      | 1         | 0          <- Reply to #1\n  3  | \"Agreed\"     | 42      | 1         | 0          <- Reply to #1\n\nTop-level: findByPostIdAndParentIsNull()\nReplies:   findByParentIdOrderByCreatedAtAsc()\nDelete:    Recursive - delete all replies first", 11)

# === SLIDE 13: SHARE/REPOST ===
sl = content_slide("Feature: Share / Repost", C_PURPLE)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User clicks Share button","sharePost(post)"),
    ("2","POST /api/posts/{id}/share","interactionService.sharePost()"),
    ("3","InteractionService.sharePost()","@Transactional"),
    ("4","Check: already shared?","BadRequestException if duplicate"),
    ("5","original.shareCount++","Increment share counter"),
    ("6","Create REPOST Post","postType=REPOST, originalPost=original"),
    ("7","postRepository.save(repost)","INSERT INTO posts"),
    ("8","shareRepository.save(share)","Track who shared what"),
    ("9","notificationService.notifyShare()","Notify original author"),
], C_PURPLE)
cd = rect(sl, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.5), BG_CARD, C_PURPLE, Pt(2))
tf = cd.text_frame; tf.word_wrap = True
stxt(tf, "  REPOST Design", sz=14, c=C_PURPLE, b=True)
apara(tf, "  A repost is a NEW Post entity with postType=REPOST and originalPost FK pointing to the original.", sz=12, c=LGRAY)
apara(tf, "  Frontend renders reposts differently: shows original content + 'Shared by @username' header.", sz=12, c=LGRAY)

# === SLIDE 14: CONNECTIONS ===
sl = content_slide("Feature: Connection / Network System", C_BLUE)
code_blk(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(2.2),
"CONNECTION STATUS STATE MACHINE:\n\n  [Follow Request] -> PENDING ---[Accept]---> ACCEPTED (in followers/following)\n                         |------[Reject]---> REJECTED\n                         |------[Block]----> BLOCKED (all interactions blocked)\n\n  UNIQUE CONSTRAINT: (follower_id, following_id)  |  PUBLIC = instant ACCEPTED  |  PRIVATE = PENDING first", 12)
flow_table(sl, Inches(0.4), Inches(3.5), Inches(12.5), [
    ("1","User clicks Follow","followUser(userId)"),
    ("2","POST /api/users/{id}/follow","ConnectionController"),
    ("3","ConnectionService.followUser()","Check not already following"),
    ("4","Check target.privacy","PUBLIC->ACCEPTED | PRIVATE->PENDING"),
    ("5","connectionRepository.save()","INSERT INTO connections"),
    ("6","notificationService.notifyFollow()","NEW_FOLLOWER or CONNECTION_REQUEST"),
    ("7","UI: Follow -> Following/Requested","Button state change"),
], C_BLUE)

# === SLIDE 15: BOOKMARKS + NOTIFICATIONS ===
sl = content_slide("Feature: Bookmarks & Notifications", C_ORANGE)
c1 = rect(sl, Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.5), BG_CARD, C_ORANGE, Pt(1.5))
tb1 = tbox(sl, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.3)); tf1 = tb1.text_frame
stxt(tf1, "BOOKMARKS", sz=18, c=C_ORANGE, b=True)
apara(tf1, "Entity: id, user_id(FK), post_id(FK), createdAt", sz=11, c=LGRAY, fn="Consolas")
apara(tf1, "UNIQUE(user_id, post_id)", sz=11, c=C_ORANGE, fn="Consolas")
apara(tf1, "", sz=4, c=BG_CARD)
apara(tf1, "Toggle Bookmark Flow:", sz=14, c=WHITE, b=True)
for t in ["1. Click bookmark icon","2. Not saved: POST /api/bookmarks/posts/{id}","3. Saved: DELETE /api/bookmarks/posts/{id}","4. BookmarkService creates/removes entity","5. Icon toggles filled/empty","","View: GET /api/bookmarks (paginated, private)"]:
    apara(tf1, f"  {t}", sz=11, c=C_BLUE if t.startswith(("1","2","3","4","5")) else LGRAY)

c2 = rect(sl, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.5), BG_CARD, C_RED, Pt(1.5))
tb2 = tbox(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3)); tf2 = tb2.text_frame
stxt(tf2, "NOTIFICATION SYSTEM", sz=18, c=C_RED, b=True)
apara(tf2, "Entity: user(recipient), actor(trigger),", sz=11, c=LGRAY, fn="Consolas")
apara(tf2, "type, message, referenceId, isRead", sz=11, c=LGRAY, fn="Consolas")
apara(tf2, "", sz=4, c=BG_CARD)
apara(tf2, "Notification Triggers:", sz=14, c=WHITE, b=True)
for t in ["LIKE - someone likes your post","COMMENT - someone comments","SHARE - someone shares your post","NEW_FOLLOWER - someone follows you","CONNECTION_REQUEST - follow request","CONNECTION_ACCEPTED - request accepted"]:
    apara(tf2, f"  {t}", sz=11, c=C_RED)
apara(tf2, "", sz=4, c=BG_CARD)
apara(tf2, "referenceId = polymorphic FK (post/user)", sz=11, c=YELLOW)
apara(tf2, "Bell badge: polls /unread/count", sz=11, c=LGRAY)

# === SLIDE 16: USER PROFILE ===
sl = content_slide("Feature: User Profile Management", C_BLUE)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User clicks Edit Profile","Opens edit form"),
    ("2","Modifies: name, bio, picture, cover","Form fields"),
    ("3","Clicks Save","PUT /api/users/me"),
    ("4","UserService.updateProfile()","Update non-null fields only"),
    ("5","userRepository.save(user)","UPDATE users SET bio=? WHERE id=?"),
    ("6","UserMapper.toResponse(user)","Entity -> DTO (exclude password)"),
    ("7","UI refreshes profile","New data displayed"),
], C_BLUE)
code_blk(sl, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.8),
"USER ENTITY FIELDS:\n\nIdentity:    id, username, email, password (BCrypt), name\nProfile:     bio (TEXT), profilePicture, coverPhoto, location, website\nSettings:    userType (PERSONAL|CREATOR|BUSINESS), privacy (PUBLIC|PRIVATE)\nFlags:       isVerified, isActive, emailVerified\nBusiness:    businessName, category, industry, contactInfo, businessAddress\nTimestamps:  createdAt, updatedAt\n\nPrivacy: PUBLIC = anyone can follow + see posts\n         PRIVATE = follow requests + posts visible only to followers", 11)

# === SLIDE 17: MESSAGING ===
sl = content_slide("Feature: Direct Messaging", C_PURPLE)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User types message, clicks Send","sendMessage()"),
    ("2","POST /api/messages/conversations/{id}/messages","{content}"),
    ("3","MessageService.sendMessage()","Get user, find conversation"),
    ("4","Message.builder().sender().receiver()","Create Message entity"),
    ("5","messageRepository.save(message)","INSERT INTO messages"),
    ("6","Message appears in chat","Scrolls to bottom"),
], C_PURPLE)
code_blk(sl, Inches(0.4), Inches(3.6), Inches(12.5), Inches(3.2),
"KEY MESSAGING ENDPOINTS:\n\nGET  /api/messages/conversations              Get all conversations\nPOST /api/messages/conversations              Create new conversation\nGET  /api/messages/conversations/{id}/messages Get messages\nPOST /api/messages/conversations/{id}/messages Send message\nPUT  /api/messages/conversations/{id}/read     Mark as read (isRead=true)\nGET  /api/messages/unread/count                Unread count (badge)\nDELETE /api/messages/{id}                     Soft delete (isDeleted=true)\nPUT  /api/messages/{id}                        Edit message\nPOST /api/messages/{id}/react                  React to message\nPOST /api/messages/conversations/{id}/mute     Mute conversation\n\nREAD RECEIPTS: Open conversation -> PUT .../read -> all from other user = isRead=true", 11)

# === SLIDE 18: STORIES ===
sl = content_slide("Feature: Stories (24h Ephemeral Content)", C_ORANGE)
code_blk(sl, Inches(0.4), Inches(1.1), Inches(6.0), Inches(2.5),
"STORY ENTITY:\n\nid, user(FK), mediaUrl, caption (280 chars)\ncreatedAt, expiresAt (= createdAt + 24 HOURS)\nisHighlight (boolean) <- permanent if true\nviewCount (int)\n\n@PrePersist: expiresAt = now().plusHours(24)\nisExpired(): now().isAfter(expiresAt)", 11)
code_blk(sl, Inches(6.7), Inches(1.1), Inches(6.3), Inches(2.5),
"KEY ENDPOINTS:\n\nPOST /api/stories            Create story\nGET  /api/stories/feed        Followed users' stories\nGET  /api/stories/my           My stories\nDELETE /api/stories/{id}      Delete story\nPOST /api/stories/{id}/view   Mark viewed (+viewCount)\nPOST /api/stories/{id}/react  React to story\nPOST /api/stories/{id}/reply  Reply (sends DM)\nPUT  /api/stories/{id}/highlight Toggle highlight\nGET  /api/stories/highlights   Permanent stories", 11)
flow_table(sl, Inches(0.4), Inches(3.9), Inches(12.5), [
    ("1","User clicks + (Add Story)","Opens creation UI"),
    ("2","Upload image, add caption, Share","createStory()"),
    ("3","POST /api/stories {mediaUrl, caption}","StoryController"),
    ("4","StoryService.createStory()","Build Story entity"),
    ("5","storyRepository.save(story)","@PrePersist sets expiresAt"),
    ("6","Story appears in carousel","Feed filters expired stories"),
], C_ORANGE)

# === SLIDE 19: ANALYTICS ===
sl = content_slide("Feature: Analytics Dashboard", C_GREEN)
eps = [
    ("GET /api/analytics/overview","Total posts, likes, comments, shares, followers"),
    ("GET /api/analytics/profile-views?days=7","Daily profile view trends"),
    ("GET /api/analytics/post-performance?days=7","Likes, comments, shares per day"),
    ("GET /api/analytics/followers/growth?days=30","Daily follower count changes"),
    ("GET /api/analytics/engagement?days=7","Engagement rate calculation"),
    ("GET /api/analytics/audience","Audience demographics"),
    ("GET /api/analytics/best-time","Hour-by-hour engagement analysis"),
    ("GET /api/analytics/top-posts?limit=10","Ranked by (likes+comments+shares)"),
    ("GET /api/analytics/hashtag-performance","Hashtag engagement metrics"),
    ("GET /api/analytics/content-type","Performance by TEXT/IMAGE/VIDEO"),
    ("GET /api/analytics/export?format=csv","Export data as CSV/JSON"),
]
y = Inches(1.0)
for ep, desc in eps:
    s = rect(sl, Inches(0.4), y, Inches(12.5), Inches(0.48), BG_CARD)
    tf = s.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {ep}"; r1.font.size = Pt(10); r1.font.color.rgb = C_GREEN; r1.font.bold = True; r1.font.name = "Consolas"
    r2 = p.add_run(); r2.text = f"    {desc}"; r2.font.size = Pt(10); r2.font.color.rgb = LGRAY; r2.font.name = "Segoe UI"
    y += Inches(0.5)
tb = tbox(sl, Inches(0.4), Inches(6.6), Inches(12), Inches(0.4))
stxt(tb.text_frame, "Computed via aggregate queries: SUM(like_count), COUNT(posts), GROUP BY date/type", sz=11, c=YELLOW, b=True)

# === SLIDE 20: SEARCH ===
sl = content_slide("Feature: Search & Discovery", C_BLUE)
flow_table(sl, Inches(0.4), Inches(1.1), Inches(12.5), [
    ("1","User types in search bar","Debounced input"),
    ("2","GET /api/search/suggestions?query=q","Auto-complete"),
    ("3","User presses Enter","GET /api/search/all?query=q"),
    ("4","SearchService.searchAll()","Search users AND posts"),
    ("5","userRepository.searchPublicUsers(q)","LIKE %query% on username/name"),
    ("6","postRepository.findByContentContaining(q)","Content search"),
    ("7","Combined results returned","{users: [...], posts: [...]}"),
], C_BLUE)
code_blk(sl, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.8),
"ADVANCED SEARCH FILTERS:\n\nPosts:  GET /api/search/posts/advanced\n  ?query=...&author=...&dateFrom=...&dateTo=...&postType=IMAGE&minLikes=10\n\nUsers:  GET /api/search/users/advanced\n  ?query=...&location=...&userType=CREATOR&verified=true\n\nEXPLORE / TRENDING:\n  GET /api/posts/trending       -> ORDER BY (likeCount+commentCount+shareCount) DESC\n  GET /api/search/trending      -> Most popular search queries\n  GET /api/users/suggestions    -> Users not followed, public, verified\n  GET /api/posts/hashtag/{tag}  -> Posts containing specific hashtag", 11)

# === SLIDE 21: MEDIA + BUSINESS ===
sl = content_slide("Feature: Media Upload & Business/Creator Tools", C_ORANGE)
c1 = rect(sl, Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.5), BG_CARD, C_ORANGE, Pt(1.5))
tb1 = tbox(sl, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.3)); tf1 = tb1.text_frame
stxt(tf1, "MEDIA UPLOAD", sz=18, c=C_ORANGE, b=True)
for t in ["POST /api/media/upload (multipart)","POST /api/media/upload/multiple","POST /api/media/upload/profile-picture","POST /api/media/upload/cover-photo","POST /api/media/upload/video"]:
    apara(tf1, f"  {t}", sz=11, c=LGRAY, fn="Consolas")
apara(tf1, "", sz=4, c=BG_CARD)
apara(tf1, "How it works:", sz=14, c=WHITE, b=True)
apara(tf1, "Angular: FormData + multipart/form-data", sz=11, c=C_BLUE)
apara(tf1, "Spring: @RequestParam('file') MultipartFile", sz=11, c=C_GREEN)
apara(tf1, "", sz=4, c=BG_CARD)
apara(tf1, "Profile Pic Upload:", sz=14, c=WHITE, b=True)
for t in ["1. Click profile area -> file picker","2. POST multipart -> save -> URL","3. user.profilePicture = newUrl","4. userRepository.save(user)"]:
    apara(tf1, f"  {t}", sz=11, c=LGRAY)

c2 = rect(sl, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.5), BG_CARD, C_PURPLE, Pt(1.5))
tb2 = tbox(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3)); tf2 = tb2.text_frame
stxt(tf2, "BUSINESS / CREATOR TOOLS", sz=18, c=C_PURPLE, b=True)
apara(tf2, "User Types:", sz=14, c=WHITE, b=True)
apara(tf2, "  PERSONAL | CREATOR | BUSINESS", sz=11, c=C_PURPLE)
apara(tf2, "", sz=4, c=BG_CARD)
apara(tf2, "CTA Buttons: [[CTA|Buy Now|url]]", sz=12, c=C_RED, b=True)
apara(tf2, "  Stored as markers in post content", sz=10, c=LGRAY)
apara(tf2, "  Extracted by parseMetadata() on read", sz=10, c=LGRAY)
apara(tf2, "", sz=4, c=BG_CARD)
apara(tf2, "Product Tags: [[TAGS|item1,item2]]", sz=12, c=C_RED, b=True)
apara(tf2, "Promotional: [[PROMO|PartnerName]]", sz=12, c=C_RED, b=True)
apara(tf2, "", sz=4, c=BG_CARD)
apara(tf2, "Design: Schema-less metadata", sz=14, c=YELLOW, b=True)
apara(tf2, "  No extra DB columns needed", sz=10, c=LGRAY)
apara(tf2, "  No migrations for new metadata", sz=10, c=LGRAY)
apara(tf2, "  parseMetadata() extracts on read", sz=10, c=LGRAY)
apara(tf2, "  buildContent() reconstructs on write", sz=10, c=LGRAY)

# === SLIDE 22: CROSS-CUTTING ===
sl = content_slide("Cross-Cutting: AOP, Exceptions, Response Wrappers", C_RED)
code_blk(sl, Inches(0.3), Inches(1.0), Inches(6.5), Inches(3.5),
"AOP LOGGING (LoggingAspect.java):\n\n@Aspect @Component\npublic class LoggingAspect {\n\n  @Pointcut(\"within(@Repository *) ||\n            within(@Service *) ||\n            within(@RestController *)\")\n  public void springBeanPointcut() {}\n\n  @Around(\"pointcut\")\n  public Object logAround(ProceedingJoinPoint jp) {\n    log.debug(\"Enter: {}.{}()\", cls, method);\n    Object result = jp.proceed();\n    log.debug(\"Exit: {}.{}()\", cls, method);\n    return result;\n  }\n\n  @AfterThrowing -> auto-log exceptions\n}", 10)
code_blk(sl, Inches(7.0), Inches(1.0), Inches(6.0), Inches(3.5),
"API RESPONSE WRAPPER (ApiResponse<T>):\n\n{ \"success\": true,\n  \"message\": \"Success\",\n  \"data\": { ... } }   // Generic type T\n\nIf data is null -> Map.of() (no null in JSON)\nFactory: ApiResponse.success(data)\n         ApiResponse.error(msg)\n\nPAGED RESPONSE (PagedResponse<T>):\n\n{ \"content\": [...],\n  \"pageNumber\": 0, \"pageSize\": 10,\n  \"totalElements\": 150, \"totalPages\": 15,\n  \"first\": true, \"last\": false }\n\nPagedResponse.fromEntityPage(page, mapper)", 10)
code_blk(sl, Inches(0.3), Inches(4.7), Inches(12.7), Inches(1.8),
"DEPENDENCY INJECTION: @RequiredArgsConstructor (Lombok)\n  private final PostRepository postRepository;  // Spring auto-injects\n  Fields are final (immutable), easy to mock in tests\n\nEXCEPTION HANDLING: @ControllerAdvice\n  ResourceNotFoundException -> 404 | DuplicateResourceException -> 409\n  BadRequestException -> 400       | UnauthorizedException -> 403", 11)

# === SLIDE 23: E2E FLOW ===
sl = content_slide("End-to-End Flow: Login -> Post -> Like -> Notification", C_BLUE)
code_blk(sl, Inches(0.3), Inches(0.95), Inches(12.7), Inches(6.2),
"STEP 1: LOGIN\n  [User] -> POST /api/auth/login -> BCrypt.matches() -> generateToken()\n  -> localStorage.setItem('revconnect_token', jwt) -> navigate('/feed')\n\nSTEP 2: LOAD FEED\n  [FeedPage.ngOnInit()] -> GET /api/posts/personalized?page=0\n  -> JwtFilter validates -> SecurityContext set\n  -> PostService: get followingIds -> SELECT posts WHERE user_id IN (ids)\n  -> For each post: add isLikedByCurrentUser -> Return PagedResponse\n\nSTEP 3: CREATE POST\n  [User clicks Post] -> POST /api/posts {content: \"Hello #first\"}\n  -> PostService.createPost() -> INSERT INTO posts\n  -> hashtagService.processHashtags() -> UPDATE hashtags count\n  -> New post prepended to feed\n\nSTEP 4: SOMEONE LIKES IT\n  [Jane clicks Heart] -> POST /api/posts/42/like\n  -> Check not already liked -> INSERT INTO likes\n  -> UPDATE posts SET like_count+1 -> notifyLike()\n  -> INSERT INTO notifications (type=LIKE, user=john, actor=jane)\n\nSTEP 5: NOTIFICATION RECEIVED\n  [John's navbar] -> GET /api/notifications/unread/count -> 1\n  -> Bell badge shows '1' -> Click -> GET /api/notifications\n  -> 'jane liked your post' -> Click -> navigate to post #42\n  -> PUT /api/notifications/{id}/read -> badge resets", 11)

# === SLIDE 24: API REFERENCE ===
sl = content_slide("Complete API Reference (90+ Endpoints)", C_GREEN)
api_groups = [
    ("Auth (/api/auth)","6 eps","register, login, verify-email, resend-verification, forgot-password, reset-password"),
    ("Posts (/api/posts)","19 eps","CRUD, feed, trending, personalized, hashtag, pin, schedule, CTA, product-tags"),
    ("Interactions","12 eps","like/unlike post, add/edit/delete comment, like comment, replies, share"),
    ("Connections","10 eps","follow/unfollow, followers/following, accept/reject, block, mutual, stats"),
    ("Users (/api/users)","11 eps","profile CRUD, search, privacy, block, report, suggestions"),
    ("Notifications","6 eps","get all, unread, count, mark read, mark all, delete"),
    ("Messages","13 eps","conversations, send/edit/delete, read receipts, react, mute, search, attach"),
    ("Stories","12 eps","create, feed, my, user, delete, view, viewers, react, reply, highlight, archive"),
    ("Analytics","14 eps","overview, profile-views, post-perf, growth, engagement, audience, best-time, export"),
    ("Search","9 eps","global, users, posts, advanced, suggestions, trending, recent, clear"),
    ("Media","10 eps","upload single/multiple, profile-pic, cover, video, delete, details, thumbnail"),
    ("Bookmarks","4 eps","bookmark, remove, list, check status"),
]
y = Inches(1.0)
for nm, cnt, desc in api_groups:
    s = rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.45), BG_CARD)
    tf = s.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {nm}"; r1.font.size = Pt(11); r1.font.color.rgb = C_GREEN; r1.font.bold = True; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = f"  ({cnt})  "; r2.font.size = Pt(10); r2.font.color.rgb = YELLOW; r2.font.bold = True; r2.font.name = "Segoe UI"
    r3 = p.add_run(); r3.text = desc; r3.font.size = Pt(9); r3.font.color.rgb = LGRAY; r3.font.name = "Segoe UI"
    y += Inches(0.47)
tb = tbox(sl, Inches(3), Inches(6.7), Inches(7), Inches(0.4))
stxt(tb.text_frame, "Total: 90+ REST Endpoints  |  16 Controllers  |  17+ JPA Entities", sz=14, c=YELLOW, b=True, al=PP_ALIGN.CENTER)

# === SLIDES 25-28: INTERVIEW Q&A ===
qa_slide("Interview Q&A: Architecture & Security (Q1-Q6)", [
    ("Q1: Why monolithic over microservices?",
     "MVP stage: simpler deployment (single JAR), lower latency (no inter-service calls), easier debugging. Can extract services later as needed."),
    ("Q2: Explain the layered architecture.",
     "Controller (HTTP) -> Service (business logic, @Transactional) -> Repository (JPA). Strict separation of concerns, each layer only talks to layer below."),
    ("Q3: How does JWT authentication work?",
     "Login -> BCrypt validate -> generate JWT (sub:userId, exp:+24h, HMAC-SHA256) -> localStorage -> interceptor adds Bearer header -> JwtFilter validates every request."),
    ("Q4: Why disable CSRF?",
     "CSRF exploits cookie-based auth. We use JWT in Authorization header (not cookies), so CSRF attacks are not applicable."),
    ("Q5: How are passwords stored?",
     "BCrypt with random salt per password. Never store plaintext. passwordEncoder.matches(raw, hash) on login. Salt prevents rainbow table attacks."),
    ("Q6: Spring Security filter chain?",
     "JwtAuthenticationFilter (extract+validate token) -> SecurityFilterChain (permitAll for /api/auth/**, authenticated for rest) -> Controller."),
])

qa_slide("Interview Q&A: Features Deep-Dive (Q7-Q12)", [
    ("Q7: How does the personalized feed work?",
     "Get followingUserIds from connectionRepo -> add self -> SELECT posts WHERE user_id IN (ids) ORDER BY createdAt DESC -> paginated via PagedResponse."),
    ("Q8: How does like system prevent duplicates?",
     "Two levels: (1) UNIQUE CONSTRAINT on (user_id, post_id) in DB, (2) existsByUserIdAndPostId() check in service. Denormalized likeCount incremented atomically."),
    ("Q9: How do nested comments work?",
     "Self-referencing parent_id FK. Top-level: parent IS NULL. Replies: parent = parentComment. Recursive cascade delete. replyCount denormalized."),
    ("Q10: Connection requests for private accounts?",
     "PUBLIC -> status=ACCEPTED immediately. PRIVATE -> status=PENDING. Target can accept/reject. Only ACCEPTED connections used for feed and follower counts."),
    ("Q11: How are notifications generated?",
     "Created inline in service methods. likePost() -> notifyLike(). Creates Notification entity with type, message, referenceId (polymorphic FK to post/user)."),
    ("Q12: How do stories expire?",
     "@PrePersist sets expiresAt = createdAt + 24h. isExpired() checks now > expiresAt. Feed filters expired. isHighlight=true bypasses expiration."),
])

qa_slide("Interview Q&A: Database & Frontend (Q13-Q18)", [
    ("Q13: Why denormalized counters (likeCount)?",
     "Avoid expensive COUNT(*) JOINs on every feed load. Store count on Post directly. Trade-off: complex writes but O(1) reads for feed rendering."),
    ("Q14: What does ddl-auto=update do?",
     "Hibernate auto-adds new columns/tables without dropping data. Dev-friendly. Production: use 'validate' + Flyway/Liquibase for migrations."),
    ("Q15: How does Angular HTTP Interceptor work?",
     "HttpInterceptorFn in app.config. Intercepts every request, clones with Authorization: Bearer header. Handles 401 -> clear token + redirect to login."),
    ("Q16: Why lazy loading for routes?",
     "loadComponent: () => import('...').then(m => m.Component). Only downloads code on navigation. Reduces initial bundle size, improves TTI."),
    ("Q17: Builder pattern in entities?",
     "Lombok @Builder: User.builder().username('john').build(). Fluent API, avoids telescoping constructors. @Builder.Default for default values."),
    ("Q18: Purpose of @Transactional?",
     "All DB ops in method = single atomic unit. If any fails, all roll back. E.g., deletePost() deletes comments+likes+bookmarks+post atomically."),
])

qa_slide("Interview Q&A: Scenarios & Scaling (Q19-Q24)", [
    ("Q19: Two users like same post simultaneously?",
     "UNIQUE CONSTRAINT prevents duplicates. @Transactional provides row-level locking for likeCount increment. One succeeds, other gets constraint violation."),
    ("Q20: How to add real-time notifications?",
     "WebSocketConfig + STOMP already in codebase. notifyLike() also sends via messagingTemplate.convertAndSendToUser(). Frontend subscribes to /queue/notifications."),
    ("Q21: How to scale to millions of users?",
     "DB: read replicas + Redis cache. Backend: extract microservices + Kafka. Feed: fan-out on write. Search: Elasticsearch. Media: CDN. Schedule: Quartz."),
    ("Q22: Security vulnerabilities addressed?",
     "SQL Injection: JPA params. XSS: Angular sanitization. CSRF: N/A (JWT). Passwords: BCrypt. JWT theft: 24h expiry. Escalation: ownership checks in services."),
    ("Q23: CTA/Product Tags metadata design?",
     "Markers in content: [[CTA|label|url]], [[TAGS|items]]. parseMetadata() extracts on read. Schema-less = no DB migrations for new metadata types."),
    ("Q24: What is AOP in your project?",
     "@Aspect LoggingAspect uses @Around to auto-log entry/exit for all @Service/@Repository/@RestController. @AfterThrowing catches exceptions. Zero code changes."),
])

# More Q&A
qa_slide("Interview Q&A: Deep Technical (Q25-Q30)", [
    ("Q25: How does ApiResponse<T> generic wrapper work?",
     "Generic class: {success, message, data:T}. Factory methods success(data), error(msg). If data=null, returns Map.of(). Consistent format for all endpoints."),
    ("Q26: How does PagedResponse utility work?",
     "fromEntityPage(page, mapper) converts Spring Page<Entity> to PagedResponse<DTO>. Includes pageNumber, pageSize, totalElements, totalPages, first, last."),
    ("Q27: DTO pattern purpose?",
     "Separates internal entity from API response. User entity has password hash; UserResponse DTO excludes it. Prevents data leakage, decouples API from schema."),
    ("Q28: Post scheduling design?",
     "ScheduledExecutorService.schedule(delay) with in-memory ConcurrentHashMap. Status: SCHEDULED->PUBLISHED or FAILED. Limitation: lost on restart."),
    ("Q29: How does auth guard protect routes?",
     "CanActivateFn checks localStorage.getItem('revconnect_token'). If exists, navigation proceeds. If not, redirects to /login. Applied via canActivate: [authGuard]."),
    ("Q30: How does frontend know if post is liked?",
     "Backend includes isLikedByCurrentUser: boolean in PostResponse. Checks likeRepository.existsByUserIdAndPostId() for each post. Frontend shows filled/empty heart."),
])

qa_slide("Interview Q&A: Advanced Scenarios (Q31-Q35)", [
    ("Q31: What happens when user deletes account?",
     "CascadeType.ALL on User relationships cascades deletion to posts, comments, likes, bookmarks, messages, notifications, settings. orphanRemoval=true."),
    ("Q32: How does the explore page work?",
     "Trending: GET /api/posts/trending sorted by engagement. Suggested users: not followed, public, verified. Hashtag search. Content discovery for growth."),
    ("Q33: Explain the Mapper pattern.",
     "PostMapper/UserMapper convert entities to DTOs. toResponse() maps fields, adds computed flags (isLikedByCurrentUser). Keeps service layer clean."),
    ("Q34: How is the feed optimized?",
     "Denormalized counters avoid JOINs. Pagination (PageRequest) limits data. Indexes on (user_id, created_at). Personalized = IN clause on followed IDs."),
    ("Q35: What would you improve?",
     "1) Redis caching for feeds. 2) Elasticsearch for search. 3) WebSocket for real-time. 4) Persistent scheduling (Quartz). 5) CDN for media. 6) Rate limiting."),
])

# === FINAL SLIDE: SUMMARY ===
sl = title_slide("RevConnect - Summary", "Full-Stack Social Media Application", C_GREEN)
tb = tbox(sl, Inches(1.5), Inches(4.5), Inches(10), Inches(2.5)); tf = tb.text_frame
stxt(tf, "Architecture: Monolithic 3-Layer  |  Spring Boot + Angular 19 + MySQL", sz=14, c=C_GREEN, al=PP_ALIGN.CENTER)
apara(tf, "Security: JWT + Spring Security + BCrypt  |  Interceptor + Auth Guard", sz=14, c=C_RED, al=PP_ALIGN.CENTER)
apara(tf, "90+ REST Endpoints  |  17+ JPA Entities  |  35 Interview Q&A", sz=14, c=YELLOW, al=PP_ALIGN.CENTER)
apara(tf, "Design Patterns: Builder, DTO, Repository, DI, AOP, Interceptor", sz=14, c=C_PURPLE, al=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = r"c:\Users\pavan\REVCONNECT\RevConnect_Interview_Guide.pptx"
prs.save(output_path)
print(f"\nPPT saved successfully to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
